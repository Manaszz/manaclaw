from pptx import Presentation
from pptx.util import Inches, Pt

# Create a simple PowerPoint presentation summarizing the NemoClaw report
prs = Presentation()

# Title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "NemoClaw для закрытого контура"
slide.placeholders[1].text = "Стек, безопасность, интеграции и рекомендации"

# Helper to add bullet slide
def add_bullet_slide(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    first = True
    for b in bullets:
        if first:
            tf.text = b
            first = False
        else:
            p = tf.add_paragraph()
            p.text = b
            p.level = 0

# Slides
add_bullet_slide(
    "Что такое NemoClaw",
    [
        "Открытый стек NVIDIA (Apache 2.0) поверх OpenClaw",
        "Запускает OpenClaw внутри sandbox OpenShell",
        "Фокус на приватности, безопасности и управлении агентами",
    ],
)

add_bullet_slide(
    "Архитектура",
    [
        "OpenClaw как агентный движок",
        "OpenShell как безопасный рантайм (Landlock, egress-контроль)",
        "NemoClaw CLI + blueprints для развёртывания",
        "LLM-инференс через Ollama, vLLM, NeMo/NIM и др.",
    ],
)

add_bullet_slide(
    "Безопасность и управление",
    [
        "Политики ФС и сети, перехват системных вызовов",
        "Ролевые модели доступа и аудит действий агентов",
        "Версионируемые blueprints для разных ассистентов",
    ],
)

add_bullet_slide(
    "Интеграции и сценарии",
    [
        "Наследует экосистему OpenClaw (Git, Jira, Slack и др.)",
        "Работает в on-prem, в облаке и на рабочей станции",
        "Поддерживает OpenAI-совместимый inference в закрытом контуре",
    ],
)

add_bullet_slide(
    "Что нужно занести в контур",
    [
        "Исходники/образы NemoClaw и OpenShell",
        "Контейнеры OpenClaw и LLM-сервера (Ollama/vLLM)",
        "Локальные модели и коннекторы к внутренним системам",
    ],
)

add_bullet_slide(
    "Рекомендации",
    [
        "Использовать NemoClaw как enterprise-обвязку вокруг OpenClaw",
        "Для максимальной лёгкости — рассмотреть гибрид c тонким локальным агент-хостом",
        "В закрытом контуре ориентироваться на полностью локальный inference",
    ],
)

out_path = "output/nemoclaw-briefing.pptx"
import os
os.makedirs("output", exist_ok=True)
prs.save(out_path)
print(out_path)

# Create a very simple PDF using reportlab
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas

pdf_path = "output/nemoclaw-report.pdf"

c = canvas.Canvas(pdf_path, pagesize=A4)
width, height = A4

text = c.beginText(40, height - 50)
text.setFont("Helvetica", 12)
lines = [
    "NemoClaw для закрытого контура",
    "",
    "Краткое содержание:",
    "- Открытый стек NVIDIA (Apache 2.0) поверх OpenClaw и OpenShell",
    "- Добавляет политики безопасности, приватности и управления агентами",
    "- Поддерживает OpenAI-совместимый inference (Ollama, vLLM, NeMo/NIM)",
    "- Подходит для on-prem и изолированных контуров при наличии локального inference",
    "",
    "Для детального описания архитектуры, безопасности и плана внедрения",
    "см. markdown-отчёт 'nemoclaw-detail-report'.",
]

for line in lines:
    text.textLine(line)

c.drawText(text)
c.showPage()
c.save()
print(pdf_path)